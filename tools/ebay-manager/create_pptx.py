# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(0, 51, 102)
GREEN = RGBColor(0, 176, 80)
ORANGE = RGBColor(255, 153, 0)
LIGHT_RED = RGBColor(192, 0, 0)
LIGHT_BLUE = RGBColor(173, 216, 230)
LIGHT_GRAY = RGBColor(242, 242, 242)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

def add_title_bar(slide, title):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.05), Inches(9), Inches(0))
    line.line.color.rgb = GREEN
    line.line.width = Pt(3)

def add_box(slide, left, top, width, height, text, bg_color, text_color=BLACK, font_size=14):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = DARK_BLUE
    shape.line.width = Pt(2)

    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

# Slide 1: Title
add_title_slide("eBay物販ビジネス", "統合システム設計書（全ツール・フロー完全版）")

# Slide 2: Architecture
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "システムアーキテクチャ - 3層構造")

add_box(slide, 0.5, 1.5, 9, 0.8, "レイヤー3：UI層 - app.py（Streamlit 1000+行）", LIGHT_BLUE)
add_box(slide, 0.5, 2.6, 9, 0.8, "レイヤー2：データベース層 - monitor.db（SQLite）+ .company/", LIGHT_GRAY)
add_box(slide, 0.5, 3.7, 4.3, 0.8, "定時実行5個", GREEN, WHITE)
add_box(slide, 5.2, 3.7, 4.3, 0.8, "随時実行9個", ORANGE, WHITE)
add_box(slide, 0.5, 4.8, 9, 0.8, "レイヤー1：ビジネスロジック層（41個のツール）", DARK_BLUE, WHITE)
add_box(slide, 0.5, 5.9, 9, 0.8, "レイヤー0：実行管理層 - daily_scheduler.py", LIGHT_RED, WHITE)

# Slide 3: Daily Schedule Timeline
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "定時実行スケジュール（朝5:00-6:00）")

timeline_y = 2.0
add_box(slide, 0.5, timeline_y, 1.8, 0.6, "5:00\n秘書ルーティン", GREEN, WHITE, 12)
add_box(slide, 2.5, timeline_y, 1.8, 0.6, "5:05\neBay同期", ORANGE, WHITE, 12)
add_box(slide, 4.5, timeline_y, 1.8, 0.6, "5:10\n在庫チェック", LIGHT_RED, WHITE, 12)
add_box(slide, 6.5, timeline_y, 1.8, 0.6, "6:00\n完了", GREEN, WHITE, 12)

sub_y = 2.9
add_box(slide, 0.3, sub_y, 2.2, 1.5, "1. メール確認\n2. TODO繰越\n3. リサーチ\n   (3層構造)", LIGHT_BLUE, BLACK, 11)
add_box(slide, 2.3, sub_y, 2.2, 1.5, "498件\nメトリクス同期\n\n自動ランク\n計算", LIGHT_BLUE, BLACK, 11)
add_box(slide, 4.3, sub_y, 2.2, 1.5, "348件\n仕入先\n在庫チェック\n\nCSV/JSON\n出力", LIGHT_BLUE, BLACK, 11)
add_box(slide, 6.3, sub_y, 3.2, 1.5, "ダッシュボード\n自動更新\n\nユーザー起床時\n全タスク完了", LIGHT_BLUE, BLACK, 11)

# Slide 4: Daily Research
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "デイリーリサーチ（3層構造）- 新機能")

add_box(slide, 0.5, 1.5, 2.8, 0.6, "A. 新商品候補", ORANGE, WHITE, 14)
add_box(slide, 3.6, 1.5, 2.8, 0.6, "B. AIニュース", GREEN, WHITE, 14)
add_box(slide, 6.7, 1.5, 2.8, 0.6, "C. 制約解決", LIGHT_RED, WHITE, 14)

add_box(slide, 0.5, 2.3, 2.8, 2.5, "eBay検索\n\n利益率\nフィルタリング\n\nランク計算", LIGHT_BLUE, BLACK, 12)
add_box(slide, 3.6, 2.3, 2.8, 2.5, "WebSearch\n\nClaude/MCP\nVision\nWeb API\n最新情報自動取得", LIGHT_BLUE, BLACK, 12)
add_box(slide, 6.7, 2.3, 2.8, 2.5, "制約ファイル\nチェック\n\n解決可能か\n自動判定\n\n改善案提案", LIGHT_BLUE, BLACK, 12)

result_text = "出力: notes/に詳細記録 + ai-news-summary.md + improvement-opportunities.md"
add_box(slide, 0.5, 5.1, 9, 0.8, result_text, LIGHT_GRAY, BLACK, 12)

# Slide 5: Tools Overview
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "ツール・機能インベントリ（41個）")

add_box(slide, 0.5, 1.5, 2.1, 1.2, "定時実行\n\n5個", GREEN, WHITE, 16)
add_box(slide, 2.8, 1.5, 2.1, 1.2, "随時実行\n\n9個", ORANGE, WHITE, 16)
add_box(slide, 5.1, 1.5, 2.1, 1.2, "コア\nロジック\n4個", LIGHT_BLUE, BLACK, 16)
add_box(slide, 7.4, 1.5, 2.1, 1.2, "eBay\n連携\n\n6個", LIGHT_RED, WHITE, 16)

add_box(slide, 0.5, 3.0, 9, 0.5, "定時実行タスク", DARK_BLUE, WHITE, 12)
add_box(slide, 0.5, 3.7, 9, 1.8, "秘書ルーティン | メール確認 | デイリーリサーチ | eBay同期 | 在庫チェック", LIGHT_GRAY, BLACK, 11)

add_box(slide, 0.5, 5.7, 9, 0.5, "随時実行ツール", DARK_BLUE, WHITE, 12)
add_box(slide, 0.5, 6.4, 9, 0.8, "商品リサーチ | 結果処理 | データ拡張 | 仕入先選定 | 価格計算 | アラート | ニュース | 競合検出", LIGHT_GRAY, BLACK, 10)

# Slide 6: eBay Integration
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "eBay連携システム")

add_box(slide, 0.5, 1.5, 4.3, 0.6, "eBay API連携", DARK_BLUE, WHITE, 14)
add_box(slide, 5.2, 1.5, 4.3, 0.6, "SQLiteデータベース", DARK_BLUE, WHITE, 14)

modules = ["ebay_client.py", "ebay_sync.py", "rank_calculator.py", "database.py", "notifier.py", "runner.py"]
y_pos = 2.3
for i, module in enumerate(modules):
    if i < 3:
        add_box(slide, 0.5, y_pos + (i * 0.9), 4.3, 0.7, module, LIGHT_BLUE, BLACK, 11)
    else:
        add_box(slide, 5.2, y_pos + ((i-3) * 0.9), 4.3, 0.7, module, LIGHT_BLUE, BLACK, 11)

add_box(slide, 0.5, 5.9, 9, 0.8, "monitor.db: ebay_listings(498件) + competitor_products(348件) + monitored_items", LIGHT_GRAY, BLACK, 11)

# Slide 7: Database
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "データベース設計（3つのテーブル）")

add_box(slide, 0.5, 1.5, 2.8, 0.5, "ebay_listings", ORANGE, WHITE, 12)
add_box(slide, 0.5, 2.1, 2.8, 2.0, "498件のeBay出品\n\nitem_id, sku,\ntitle, price,\nwatch_count,\nview_count,\nrank(S-E)", LIGHT_BLUE, BLACK, 10)

add_box(slide, 3.6, 1.5, 2.8, 0.5, "competitor_products", ORANGE, WHITE, 12)
add_box(slide, 3.6, 2.1, 2.8, 2.0, "348件の仕入先商品\n\nsku, source,\nstatus(在庫有無),\nprice_jpy,\nlast_checked_at", LIGHT_BLUE, BLACK, 10)

add_box(slide, 6.7, 1.5, 2.8, 0.5, "monitored_items", ORANGE, WHITE, 12)
add_box(slide, 6.7, 2.1, 2.8, 2.0, "ユーザー監視中の商品\n\nalert設定\nアラート時刻", LIGHT_BLUE, BLACK, 10)

add_box(slide, 0.5, 4.5, 9, 0.5, "データフロー", DARK_BLUE, WHITE, 12)

flow_items = ["朝5:05\neBay同期", "メトリクス\n更新", "ランク\n計算", "朝5:10\n在庫チェック", "結果\n記録", "ダッシュボード"]
flow_x = 0.5
for item in flow_items:
    add_box(slide, flow_x, 5.2, 1.4, 1.2, item, LIGHT_GRAY if "ダッシュ" not in item else GREEN, BLACK, 10)
    flow_x += 1.5

# Slide 8: UI Dashboard
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "UI・ダッシュボード（8つのページ）")

pages = [
    ("ページ1\nダッシュボード", DARK_BLUE),
    ("ページ2\neBay連携", GREEN),
    ("ページ3\n利益管理", ORANGE),
    ("ページ4\nリサーチ", LIGHT_RED),
    ("ページ5\nSKU管理", LIGHT_BLUE),
    ("ページ6\n競合分析", DARK_BLUE),
    ("ページ7\n学習管理", GREEN),
    ("ページ8\nログ", ORANGE)
]

page_x = 0.5
page_y = 1.8
for i, (page_name, color) in enumerate(pages):
    if i == 4:
        page_x = 0.5
        page_y = 4.5

    text_color = WHITE if color != LIGHT_BLUE else BLACK
    add_box(slide, page_x, page_y, 1.8, 1.5, page_name, color, text_color, 12)
    page_x += 2.0

# Slide 9: Error Handling
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "エラーハンドリング（3段階）")

add_box(slide, 0.5, 1.5, 2.8, 0.6, "レベル1", GREEN, WHITE, 14)
add_box(slide, 0.5, 2.2, 2.8, 1.8, "自動リカバリー\n\nTimeout → リトライ\nAPI制限 → 待機", LIGHT_BLUE, BLACK, 11)

add_box(slide, 3.6, 1.5, 2.8, 0.6, "レベル2", ORANGE, WHITE, 14)
add_box(slide, 3.6, 2.2, 2.8, 1.8, "ログ記録\n部分継続\n\n単一失敗 → 続行\nページなし → 記録", LIGHT_BLUE, BLACK, 11)

add_box(slide, 6.7, 1.5, 2.8, 0.6, "レベル3", LIGHT_RED, WHITE, 14)
add_box(slide, 6.7, 2.2, 2.8, 1.8, "手動介入必要\n\n認証失敗 → メール\nクラッシュ → 停止", LIGHT_BLUE, BLACK, 11)

add_box(slide, 0.5, 4.3, 9, 0.5, "影響度の比較", DARK_BLUE, WHITE, 12)

impact_y = 5.0
add_box(slide, 0.5, impact_y, 2.8, 1.0, "レベル1\n自動解決", GREEN, WHITE, 10)
add_box(slide, 3.5, impact_y, 5.8, 1.0, "対応時間: 0秒～数秒（ユーザー無意識）", LIGHT_GRAY, BLACK, 11)

impact_y = 6.1
add_box(slide, 0.5, impact_y, 2.8, 0.7, "レベル2\nユーザー確認", ORANGE, WHITE, 10)
add_box(slide, 3.5, impact_y, 5.8, 0.7, "対応時間: 数分～数時間", LIGHT_GRAY, BLACK, 11)

# Slide 10: Checklist
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "実装チェックリスト")

checklist = [
    ("monitor/scrapers.py 削除", True),
    ("monitor.db 初期化", False),
    ("daily_scheduler.py スケジュール設定", False),
    ("eBay API credentials 設定", False),
    ("Gmail OAuth 設定", False),
    ("Selenium WebDriver インストール", False),
    ("エラー通知テスト", False),
    ("app.py 動作確認", False),
]

check_y = 1.5
for item, done in checklist:
    status = "DONE" if done else "TODO"
    status_color = GREEN if done else ORANGE
    add_box(slide, 0.5, check_y, 0.8, 0.4, status, status_color, WHITE, 9)
    add_box(slide, 1.5, check_y, 8.0, 0.4, item, LIGHT_GRAY, BLACK, 11)
    check_y += 0.5

# Slide 11: Future Phases
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "拡張計画（Phase 5-8）")

phases = [
    ("Phase 5\nリアルタイム通知", "Discord/Slack/メール", GREEN),
    ("Phase 6\nAI推奨システム", "買い時・売り時予測", ORANGE),
    ("Phase 7\nマルチマーケット", "Amazon/Mercado Libre", LIGHT_RED),
    ("Phase 8\n完全自動化", "自動価格・出品・仕入れ", DARK_BLUE)
]

phase_x = 0.5
for phase_name, description, color in phases:
    text_color = WHITE if color != ORANGE else BLACK

    add_box(slide, phase_x, 1.8, 2.1, 0.8, phase_name, color, text_color, 11)
    add_box(slide, phase_x, 2.7, 2.1, 2.8, description, LIGHT_BLUE, BLACK, 10)
    phase_x += 2.3

# Slide 12: Summary
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "システム統合の成果")

summary_items = [
    ("41個のツール", "すべてが統合されたシステム"),
    ("朝5:00-6:00", "完全自動実行で全タスク完了"),
    ("随時実行9個", "強力な分析・管理ツール"),
    ("8ページUI", "全機能をワンストップで操作"),
    ("SQLite管理", "堅牢なデータベース設計"),
]

summary_y = 1.8
for title, description in summary_items:
    add_box(slide, 0.5, summary_y, 2.2, 0.7, title, DARK_BLUE, WHITE, 11)
    add_box(slide, 2.9, summary_y, 6.6, 0.7, description, LIGHT_BLUE, BLACK, 11)
    summary_y += 0.9

prs.save("ebay-manager-system-design.pptx")
print("[OK] 図表付きPowerPoint資料を日本語で作成しました")
print("[FILE] ebay-manager-system-design.pptx")
print("[SLIDES] 12スライド - ビジュアル重視の設計")
